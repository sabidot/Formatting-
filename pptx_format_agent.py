"""
LightCastle PPT Formatting Agent
=================================
Rules (from LightCastle_Deck_Template_2025.pptx):

FONT
  - Always Lato or Lato Black — never changed

SIZE (title / subtitle only)
  - Find the majority size for each role group (title = Lato Black, subtitle = coloured header text)
  - Fix minority outliers to match majority
  - Never change body text sizes (author's choice for fitting text boxes)

COLOR
  - Reset non-standard colors to nearest design color
  - Keep all LC brand palette colors intact
  - Keep bold / italic / highlight intact

TABLES
  - Header row    : center + middle
  - Numeric cells : right  + middle
  - All others    : left   + middle
  - Decimal numbers in cells → round to nearest integer

FLAG LOCATIONS
  - Always show "Slide X" (not Para/Table coords)
"""

import re
import json
import sys
import argparse
from pathlib import Path
from collections import Counter, defaultdict

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ---------------------------------------------------------------------------
# SPEC
# ---------------------------------------------------------------------------

# LC brand palette — NEVER reset these
DESIGN_COLORS = {
    '0473BA', '429DE5', '05519B', '2E75B6', '0D3395',
    '055291', '020964', '061B39', '002060', '002350',
    '1C4587', 'FBBF24', '737373', '595959', 'FFFFFF',
    '000000', '1E1E1E', 'EEF5FD', 'F2F2F2', 'F9F9F9',
    'DAF8FF', '068ADC', '00B0F0', '00172D', '212121',
    '262626', '3F3F3F', 'AEAEAF', '78909C', '4285F4',
    'FFAB40',
}

# Standard body text colors
BODY_COLORS = {'000000', '1E1E1E'}

# Colors that are clearly wrong (not brand, not near-black)
# anything not in DESIGN_COLORS gets flagged / reset

SEVERITY_ORDER = {"critical": 0, "high": 1, "medium": 2, "low": 3}

NUMERIC_RE = re.compile(r'^\s*[\d,\.\s%$\-\+]+\s*$')
DECIMAL_RE  = re.compile(r'\b\d+\.\d+\b')


# ---------------------------------------------------------------------------
# HELPERS
# ---------------------------------------------------------------------------

def _safe_slides(prs):
    for i in range(len(prs.slides)):
        try:
            yield i, prs.slides[i]
        except Exception:
            continue


def _get_color(run):
    try:
        if run.font.color and run.font.color.type:
            return str(run.font.color.rgb).upper()
    except Exception:
        pass
    return None


def _get_size_pt(run):
    try:
        if run.font.size:
            return round(run.font.size.pt, 2)
    except Exception:
        pass
    return None


def _is_copyright(run):
    try:
        name = run.font.name or ''
        sz   = _get_size_pt(run)
        text = run.text
        return (
            'Poppins' in name or
            'LightCastle Partners' in text or
            '©' in text or 'Ⓒ' in text or
            (sz and sz <= 8 and _get_color(run) == 'FFFFFF')
        )
    except Exception:
        return False


def _is_title_run(run):
    name = run.font.name or ''
    sz   = _get_size_pt(run)
    return 'Black' in name and sz and sz >= 24


def _is_subtitle_run(run):
    sz    = _get_size_pt(run)
    color = _get_color(run)
    name  = run.font.name or ''
    return (
        sz and 14 <= sz <= 24 and
        'Black' not in name and
        color in ('737373', '429DE5', 'FFFFFF', None)
    )


# ---------------------------------------------------------------------------
# MAJORITY-SIZE DETECTION
# ---------------------------------------------------------------------------

def _collect_role_sizes(prs):
    """
    Collect all (slide, run) pairs grouped by role fingerprint.
    Returns dict: role_key -> [(slide_num, run_size)]
    Role key = ('title'|'subtitle', color_or_none)
    """
    groups = defaultdict(list)

    for s_idx, slide in _safe_slides(prs):
        for shape in slide.shapes:
            try:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if not run.text.strip() or _is_copyright(run):
                            continue
                        sz    = _get_size_pt(run)
                        color = _get_color(run)
                        if sz is None:
                            continue
                        if _is_title_run(run):
                            if color not in _UNIFORM_SKIP_COLORS:
                                groups[('title', color)].append((s_idx + 1, sz, run))
                        elif _is_subtitle_run(run):
                            if color not in _UNIFORM_SKIP_COLORS:
                                groups[('subtitle', color)].append((s_idx + 1, sz, run))
            except Exception:
                continue

    return groups


def _majority_sizes(groups):
    """
    For each group, compute the majority size.
    Returns dict: role_key -> majority_size_pt
    """
    result = {}
    for key, entries in groups.items():
        sizes = [sz for _, sz, _ in entries]
        if sizes:
            result[key] = Counter(sizes).most_common(1)[0][0]
    return result


# ---------------------------------------------------------------------------
# FLAGGING
# ---------------------------------------------------------------------------

def flag_size_drift(prs, flags):
    """
    For title and subtitle runs, flag any that don't match the majority size
    for their role group.
    """
    groups   = _collect_role_sizes(prs)
    majority = _majority_sizes(groups)

    for key, entries in groups.items():
        role, color = key
        maj_sz = majority.get(key)
        if maj_sz is None:
            continue
        for (slide_num, sz, run) in entries:
            if sz != maj_sz:
                flags.append({
                    "type":     f"{role}_size_drift",
                    "severity": "high",
                    "slide":    slide_num,
                    "found":    f"{sz}pt",
                    "expected": f"{maj_sz}pt (majority)",
                    "preview":  run.text[:50],
                    "fix":      f"Set to {maj_sz}pt",
                })


def flag_colors(prs, flags):
    """Flag runs with colors not in the design palette."""
    for s_idx, slide in _safe_slides(prs):
        for shape in slide.shapes:
            try:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if not run.text.strip() or _is_copyright(run):
                            continue
                        color = _get_color(run)
                        if color and color not in DESIGN_COLORS:
                            flags.append({
                                "type":     "non_standard_color",
                                "severity": "medium",
                                "slide":    s_idx + 1,
                                "found":    f"#{color}",
                                "expected": "LC brand palette",
                                "preview":  run.text[:50],
                                "fix":      "Reset to nearest standard color",
                            })
            except Exception:
                continue


def flag_table_alignment(prs, flags):
    """Flag table cells with wrong alignment."""
    for s_idx, slide in _safe_slides(prs):
        for shape in slide.shapes:
            try:
                if not shape.has_table:
                    continue
                tbl = shape.table
                for r_idx, row in enumerate(tbl.rows):
                    for c_idx, cell in enumerate(row.cells):
                        text = cell.text.strip()
                        is_header = (r_idx == 0)
                        is_num    = bool(NUMERIC_RE.match(text)) if text else False

                        exp_h = PP_ALIGN.CENTER if is_header else (PP_ALIGN.RIGHT if is_num else PP_ALIGN.LEFT)
                        exp_v = 3  # MIDDLE

                        actual_h = cell.text_frame.paragraphs[0].alignment if cell.text_frame.paragraphs else None
                        actual_v = cell.vertical_anchor

                        exp_h_name = 'center' if is_header else ('right' if is_num else 'left')

                        if actual_h != exp_h:
                            flags.append({
                                "type":     "table_h_align",
                                "severity": "high",
                                "slide":    s_idx + 1,
                                "found":    str(actual_h),
                                "expected": exp_h_name,
                                "preview":  text[:40],
                                "fix":      f"Set horizontal align to {exp_h_name}",
                            })
            except Exception:
                continue


def flag_table_decimals(prs, flags):
    """Flag decimal numbers in table cells."""
    for s_idx, slide in _safe_slides(prs):
        for shape in slide.shapes:
            try:
                if not shape.has_table:
                    continue
                tbl = shape.table
                for r_idx, row in enumerate(tbl.rows):
                    for c_idx, cell in enumerate(row.cells):
                        text = cell.text.strip()
                        if DECIMAL_RE.search(text):
                            flags.append({
                                "type":     "table_decimal",
                                "severity": "medium",
                                "slide":    s_idx + 1,
                                "found":    text[:40],
                                "expected": "Integer (no decimals)",
                                "fix":      "Round to nearest integer",
                            })
            except Exception:
                continue


def run_all_checks(prs):
    flags = []
    flag_size_drift(prs, flags)
    flag_colors(prs, flags)
    flag_table_alignment(prs, flags)
    flag_table_decimals(prs, flags)
    return flags


# ---------------------------------------------------------------------------
# AUTO-FIXER
# ---------------------------------------------------------------------------

def fix_size_drift(prs):
    """Fix title and subtitle runs to match their group majority size."""
    groups   = _collect_role_sizes(prs)
    majority = _majority_sizes(groups)

    for key, entries in groups.items():
        maj_sz = majority.get(key)
        if maj_sz is None:
            continue
        for (slide_num, sz, run) in entries:
            if sz != maj_sz:
                run.font.size = Pt(maj_sz)


def fix_colors(prs):
    """
    Reset non-palette colors to nearest standard.
    Bold / italic / highlight preserved.
    Near-white off-whites → FFFFFF.
    Near-black off-blacks → 1E1E1E.
    """
    def nearest_standard(color):
        # Simple heuristic: if very light → white, if very dark → near-black
        try:
            r = int(color[0:2], 16)
            g = int(color[2:4], 16)
            b = int(color[4:6], 16)
            brightness = (r + g + b) / 3
            if brightness > 200:
                return 'FFFFFF'
            elif brightness < 80:
                return '1E1E1E'
            else:
                return '1E1E1E'  # default to dark text
        except Exception:
            return '1E1E1E'

    for s_idx, slide in _safe_slides(prs):
        for shape in slide.shapes:
            try:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if not run.text.strip() or _is_copyright(run):
                            continue
                        color = _get_color(run)
                        if color and color not in DESIGN_COLORS:
                            std = nearest_standard(color)
                            r = int(std[0:2], 16)
                            g = int(std[2:4], 16)
                            b = int(std[4:6], 16)
                            run.font.color.rgb = RGBColor(r, g, b)
            except Exception:
                continue


def fix_table_alignment(prs):
    """Set correct alignment on all table cells."""
    for s_idx, slide in _safe_slides(prs):
        for shape in slide.shapes:
            try:
                if not shape.has_table:
                    continue
                tbl = shape.table
                for r_idx, row in enumerate(tbl.rows):
                    for c_idx, cell in enumerate(row.cells):
                        text = cell.text.strip()
                        is_header = (r_idx == 0)
                        is_num    = bool(NUMERIC_RE.match(text)) if text else False

                        # Horizontal alignment
                        if is_header:
                            h_align = PP_ALIGN.CENTER
                        elif is_num:
                            h_align = PP_ALIGN.RIGHT
                        else:
                            h_align = PP_ALIGN.LEFT

                        for para in cell.text_frame.paragraphs:
                            para.alignment = h_align

                        # Vertical alignment → MIDDLE (3)
                        from pptx.oxml.ns import qn
                        tc = cell._tc
                        tcPr = tc.find(qn('a:tcPr'))
                        if tcPr is None:
                            tcPr = etree.SubElement(tc, qn('a:tcPr'))
                        tcPr.set('anchor', 'ctr')  # ctr = middle
            except Exception:
                continue


def fix_table_decimals(prs):
    """Round decimal numbers in table cells to integers."""
    def round_decimals(text):
        def replacer(m):
            try:
                return str(round(float(m.group())))
            except Exception:
                return m.group()
        return DECIMAL_RE.sub(replacer, text)

    for s_idx, slide in _safe_slides(prs):
        for shape in slide.shapes:
            try:
                if not shape.has_table:
                    continue
                tbl = shape.table
                for r_idx, row in enumerate(tbl.rows):
                    for c_idx, cell in enumerate(row.cells):
                        text = cell.text.strip()
                        if not DECIMAL_RE.search(text):
                            continue
                        # Replace decimal text run by run
                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                if DECIMAL_RE.search(run.text):
                                    run.text = round_decimals(run.text)
            except Exception:
                continue


def apply_all_fixes_pptx(prs):
    fix_size_drift(prs)
    fix_colors(prs)
    fix_table_alignment(prs)
    fix_table_decimals(prs)


# ---------------------------------------------------------------------------
# REPORT + CLI
# ---------------------------------------------------------------------------

def build_report(filename, flags):
    return {
        "file":       filename,
        "total_flags": len(flags),
        "severity_summary": {
            sev: sum(1 for f in flags if f.get("severity") == sev)
            for sev in ("critical", "high", "medium", "low")
        },
        "flags": sorted(flags, key=lambda f: SEVERITY_ORDER.get(f.get("severity", "low"), 9)),
    }


def main():
    parser = argparse.ArgumentParser(description="LightCastle PPT Formatting Agent")
    parser.add_argument("input",     help="Path to .pptx file")
    parser.add_argument("--check",   action="store_true", help="Flag only, no fix")
    parser.add_argument("--verbose", action="store_true", help="Print flags to console")
    args = parser.parse_args()

    path = Path(args.input)
    if not path.exists():
        print(f"Error: not found: {path}", file=sys.stderr); sys.exit(1)

    print(f"Loading: {path}")
    prs = Presentation(str(path))
    print(f"  Slides: {len(prs.slides)}")

    print("Scanning...")
    flags = run_all_checks(prs)
    print(f"  Flags: {len(flags)}")

    report = build_report(str(path), flags)
    json_path = path.with_name(path.stem + "_flags.json")
    json_path.write_text(json.dumps(report, indent=2, ensure_ascii=False), encoding="utf-8")
    print(f"  Report: {json_path}")

    if not args.check:
        print("Applying fixes...")
        apply_all_fixes_pptx(prs)
        out_path = path.with_name(path.stem + "_fixed.pptx")
        prs.save(str(out_path))
        print(f"  Fixed: {out_path}")

    if args.verbose:
        print()
        if not flags:
            print("No issues found.")
        else:
            print(f"{'SEV':<10} {'TYPE':<30} {'SLIDE':<8} DETAIL")
            print("-" * 80)
            for f in report["flags"]:
                sev   = f.get("severity","?").upper()[:8]
                ftype = f.get("type","")[:29]
                slide = str(f.get("slide","?"))
                parts = []
                if "found"   in f: parts.append(f"found={f['found']}")
                if "expected"in f: parts.append(f"expected={f['expected']}")
                if "preview" in f: parts.append(f'"{f["preview"]}"')
                if "message" in f: parts.append(f["message"])
                print(f"{sev:<10} {ftype:<30} {slide:<8} {' | '.join(parts)}")

    print()
    print("Done.")
    for sev, cnt in report["severity_summary"].items():
        if cnt: print(f"  {sev}: {cnt}")


if __name__ == "__main__":
    main()
